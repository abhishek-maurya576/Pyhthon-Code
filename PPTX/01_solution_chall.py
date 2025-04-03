from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Accessible Education for Underserved Communities"
subtitle.text = "A Technology-Driven Solution for Quality Education\nPresented by: Knights"

# Slide 2: Problem Statement
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Problem Statement"
content.text = ("• Lack of Access to Quality Education\n"
                        "• Millions of children and adults in underserved communities lack access to quality education due to:\n"
                        "  - Inadequate infrastructure\n  - Insufficient resources\n  - Shortage of trained teachers\n"
                        "• Impact: Educational gap leading to poverty and inequality")

# Slide 3: Objective
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Objective"
content.text = ("• Develop an innovative, scalable, and technology-driven solution\n"
                        "• Ensure education is accessible, inclusive, and personalized for diverse learners\n"
                        "• Align with UN SDG 4: Quality Education to promote lifelong learning opportunities for all")

# Slide 4: Our Solution
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Our Solution"
content.text = ("• A mobile-based learning platform designed for underserved communities\n"
                        "• Key Features:\n"
                        "  - Interactive lessons via Jetpack Compose\n"
                        "  - Live & recorded sessions powered by Firebase\n"
                        "  - Offline access using SQLite\n"
                        "  - AI-driven personalized learning paths")

# Slide 5: Technology Stack
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Technology Stack"
content.text = ("• Frontend: Kotlin, Jetpack Compose\n"
                        "• Backend: Firebase, SQLite\n"
                        "• Database: Firebase Firestore & SQLite\n"
                        "• Deployment: Android-based mobile application")

# Slide 6: Features
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Features"
content.text = ("• Smart Learning Modules: AI-powered personalized lessons\n"
                        "• Offline Learning: Download lessons for later use\n"
                        "• Community Forum: Peer-to-peer discussion and teacher support\n"
                        "• Progress Tracking: Personalized performance reports")

# Slide 7: MVP and Prototype
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "MVP and Prototype"
content.text = ("• MVP development is in progress\n"
                        "• Tools Used: Kotlin, Jetpack Compose, Firebase\n"
                        "• Prototype Available on Figma: https://www.figma.com/proto/Cin2qEByuqXQizsK8JfNzv/EduReach")

# Slide 8: Expected Impact
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Expected Impact"
content.text = ("• Bridging the Education Gap: Reaching underserved communities\n"
                        "• Empowering Students: Equipping learners with essential skills\n"
                        "• Scalability: Potential expansion across multiple regions")

# Slide 9: Team & Responsibilities
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Team & Responsibilities"
content.text = ("• Abhishek Maurya - App Development (Jetpack Compose, Firebase)\n"
                        "• Kumar Manglam - Backend & Database (SQLite, Firestore)\n"
                        "• Shivank Rastogi - UI/UX Design (Figma, User Research)")

# Slide 10: Conclusion & Next Steps
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion & Next Steps"
content.text = ("• Further Development: Refining features based on user feedback\n"
                        "• Testing & Deployment: Pilot launch in select regions\n"
                        "• Seeking Support: Collaboration with NGOs & Education Institutions")

# Slide 11: Thank You Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Thank You!"
content.text = ("• Contact: maurya972137@gmail.com\n"
                        "• GitHub Repository: https://github.com/abhishek-maurya576/EduReach\n"
                        "• Demo: https://www.figma.com/proto/Cin2qEByuqXQizsK8JfNzv/EduReach")

# Save the presentation
pptx_path = "Solution_Challenge.pptx"
prs.save(pptx_path)

print(f"Presentation saved to: {pptx_path}")